import streamlit as st
import docx
import json
import requests
import time
import os
import io
import urllib.parse
import random
from dotenv import load_dotenv
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# ==========================================
# 1. CẤU HÌNH HỆ THỐNG
# ==========================================
st.set_page_config(page_title="OCBS AI Report", layout="wide")

load_dotenv()
LLAMA3_API_KEY = os.getenv("GROQ_API_KEY")
PEXELS_API_KEY = os.getenv("PEXELS_API_KEY")
LLAMA3_BASE_URL = "https://api.groq.com/openai/v1" 

THEME = {
    "bg": RGBColor(255, 255, 255),        
    "title": RGBColor(0, 51, 102),        
    "card_bg": RGBColor(248, 250, 252),   
    "card_border": RGBColor(212, 175, 55),
    "text": RGBColor(51, 65, 85),         
    "accent": RGBColor(0, 153, 117),      
    "chart": [RGBColor(0, 153, 117), RGBColor(212, 175, 55), RGBColor(0, 51, 102), RGBColor(148, 163, 184)]
}
client = OpenAI(api_key=LLAMA3_API_KEY, base_url=LLAMA3_BASE_URL)

# ==========================================
# 2. KIẾN TRÚC LẤY ẢNH VÀ LOGO
# ==========================================
def download_company_logo(domain, filepath):
    try:
        res = requests.get(f"https://logo.clearbit.com/{domain}?size=800", timeout=10)
        if res.status_code == 200:
            with open(filepath, 'wb') as f: f.write(res.content)
            return True
    except: pass
    return False

def generate_ai_image(image_prompt, filepath):
    try:
        style = "modern corporate finance, high-tech glass building, glowing digital data charts, professional lighting, cinematic, hyper-realistic, 8k resolution, clean minimalist design --ar 16:9"
        encoded = urllib.parse.quote(f"{image_prompt}, {style}")
        headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}
        res = requests.get(f"https://image.pollinations.ai/prompt/{encoded}?width=1024&height=576&nologo=true&seed={random.randint(1, 100000)}", headers=headers, timeout=20)
        if res.status_code == 200 and len(res.content) > 5000: 
            with open(filepath, 'wb') as f: f.write(res.content)
            return True
    except: pass
    return False

def download_image_pexels(keyword, filepath):
    try:
        res = requests.get(f"https://api.pexels.com/v1/search?query={keyword} finance business&per_page=10&orientation=landscape", headers={"Authorization": PEXELS_API_KEY})
        photos = res.json().get("photos", [])
        if photos:
            img_data = requests.get(random.choice(photos)["src"]["large"]).content
            with open(filepath, 'wb') as f: f.write(img_data)
            return True
    except: pass
    return False

def get_image_robust(prompt, keyword, filepath):
    if generate_ai_image(prompt, filepath): return True
    return download_image_pexels(keyword, filepath)

# ==========================================
# 3. LÕI AI: ĐỌC WORD VÀ ÉP KHUÔN JSON
# ==========================================
def get_semantic_chunks_from_docx(file_stream, max_words=600):
    file_stream.seek(0) 
    doc = docx.Document(file_stream)
    chunks, curr, count = [], [], 0
    for para in doc.paragraphs:
        txt = para.text.strip()
        if txt: 
            curr.append(txt + " ") 
            count += len(txt.split())
            if count >= max_words:
                chunks.append("\n\n".join(curr))
                curr, count = [], 0
    if curr: chunks.append("\n\n".join(curr))
    return chunks

def get_slide_json_from_llama3(file_stream, status_container):
    chunks = get_semantic_chunks_from_docx(file_stream)
    all_slides = []
    
    system_prompt = """
    Chuyển văn bản thành JSON Presentation. TUYỆT ĐỐI GIỮ TIẾNG VIỆT CÓ DẤU.
    BẮT BUỘC TRẢ VỀ JSON THEO ĐÚNG MẪU NÀY:
    {
        "slides": [
            {
                "title": "Tiêu đề Slide",
                "type": "text", 
                "company_domain": "",
                "image_prompt": "10 từ tiếng anh miêu tả",
                "takeaway": "Một câu chốt ngắn gọn",
                "positive_stocks": [],
                "negative_stocks": [],
                "bullets": ["Ý 1", "Ý 2", "Ý 3"],
                "table_data": [],
                "chart_data": {}
            }
        ]
    }
    
    QUY TẮC SỐNG CÒN:
    1. PHÂN LỚP TYPE: 
       - Lộ Trình/Giai đoạn/Kế hoạch năm -> "type": "table", "table_data": [["Mốc thời gian", "Sự kiện"]] (Mốc thời gian chỉ được ghi siêu ngắn gọn, ví dụ: "Năm 2026").
       - Tỷ trọng/% -> "type": "chart", "chart_data": {"Nội địa": 60, "Quốc tế": 40}
       - Còn lại -> "type": "text" hoặc "grid"
    2. CHỐNG RỖNG: Nếu đoạn văn quá ngắn hoặc chỉ có tiêu đề, HÃY TỰ SUY LUẬN VÀ TÓM TẮT THÀNH ÍT NHẤT 2 Ý ĐIỀN VÀO 'bullets'. TUYỆT ĐỐI KHÔNG ĐỂ MẢNG BULLETS RỖNG [].
    """
    
    for i, chunk in enumerate(chunks):
        status_container.write(f"Đang bóc tách dữ liệu JSON (Gói {i+1}/{len(chunks)})...")
        retries = 3
        while retries > 0:
            try:
                res = client.chat.completions.create(
                    model="llama-3.1-8b-instant", 
                    messages=[{"role": "system", "content": system_prompt}, {"role": "user", "content": chunk}],
                    temperature=0.2, response_format={"type": "json_object"}
                )
                content = res.choices[0].message.content.strip()
                
                content = content.replace("```json", "").replace("```", "").strip()
                chunk_json = json.loads(content)
                
                if isinstance(chunk_json, dict):
                    if "slides" in chunk_json: all_slides.extend(chunk_json["slides"])
                    elif "slide" in chunk_json: all_slides.extend(chunk_json["slide"])
                    elif "presentation" in chunk_json: all_slides.extend(chunk_json["presentation"])
                    elif "data" in chunk_json: all_slides.extend(chunk_json["data"])
                    elif "title" in chunk_json: all_slides.append(chunk_json)
                elif isinstance(chunk_json, list): all_slides.extend(chunk_json)
                    
                break
            except Exception as e:
                retries -= 1
                if retries == 0: 
                    st.error(f"Lỗi JSON gói {i+1}: {e}")
                    st.code(content, language="json")
                    status_container.error(f"Lỗi bóc tách JSON ở gói {i+1}.")
                    
    return {"slides": all_slides}

# ==========================================
# 4. CÁC HÀM HELPER VẼ PPTX
# ==========================================
def clear_placeholders(slide):
    for shape in slide.placeholders:
        sp = shape.element
        sp.getparent().remove(sp)

def get_dynamic_pt(text, default_size):
    length = len(str(text))
    if length < 30: return Pt(default_size + 4)
    elif length < 80: return Pt(default_size)
    elif length < 150: return Pt(max(12, default_size - 3))
    else: return Pt(max(11, default_size - 6))

def set_p_format(paragraph, text, font_size, bold=False, color_rgb=None, alignment=None):
    paragraph.text = str(text)
    font = paragraph.font
    font.name = 'Arial'
    font.size = font_size
    font.bold = bold
    if color_rgb: font.color.rgb = color_rgb
    if alignment: paragraph.alignment = alignment
    
    if hasattr(paragraph, '_parent'):
        tf = paragraph._parent
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

def add_editable_stock_tag(slide, x, y, stock_code, trend):
    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.1), Inches(0.4))
    bg_color = THEME["accent"] if trend == 'up' else RGBColor(220, 38, 38)
    tag.fill.solid(); tag.fill.fore_color.rgb = bg_color; tag.line.fill.background = None
    
    tf = tag.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    set_p_format(tf.paragraphs[0], stock_code.upper(), get_dynamic_pt(stock_code, 12), True, RGBColor(255, 255, 255), PP_ALIGN.CENTER)

def draw_editable_ocbs_grid(slide, x, y, w, h, index, text):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid(); card.fill.fore_color.rgb = THEME["card_bg"]; card.line.color.rgb = THEME["card_border"]
    
    tf = card.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.15)
    set_p_format(tf.paragraphs[0], f"{index}. {text}", get_dynamic_pt(text, 16), False, THEME["text"])

# ==========================================
# 5. RENDER PPTX CLEAN
# ==========================================
def render_pptx_clean(slide_data, template_path, status_container):
    status_container.write("Đang Render Giao diện Chống tràn lề...")
    prs = Presentation(template_path)
    SAFE_LEFT, SAFE_TOP = Inches(0.5), Inches(0.8)
    AVAIL_W = prs.slide_width - Inches(1.0)
    
    try:
        ts = prs.slides.add_slide(prs.slide_layouts[6])
        clear_placeholders(ts) 
        set_p_format(ts.shapes.add_textbox(SAFE_LEFT, Inches(3.0), AVAIL_W, Inches(2.0)).text_frame.paragraphs[0], "CHIẾN LƯỢC CHUYỂN ĐỔI SỐ OCBS", Pt(54), True, THEME["title"], PP_ALIGN.CENTER)
    except: pass

    slides_list = slide_data.get('slides', slide_data) if isinstance(slide_data, dict) else slide_data

    for i, s_info in enumerate(slides_list):
        if not isinstance(s_info, dict): continue
        try:
            s_type = str(s_info.get('type') or 'text')
            title = str(s_info.get('title') or 'Untitled')
            takeaway = str(s_info.get('takeaway') or '')
            bullets = [str(b).strip() for b in (s_info.get('bullets') or []) if b]
            if not bullets: bullets = ["Nội dung đang được cập nhật."]
            
            p_stocks = s_info.get('positive_stocks') or []
            n_stocks = s_info.get('negative_stocks') or []

            try: slide = prs.slides.add_slide(prs.slide_layouts[6]); clear_placeholders(slide) 
            except: slide = prs.slides.add_slide(prs.slide_layouts[1]); clear_placeholders(slide)

            p_title = slide.shapes.add_textbox(SAFE_LEFT, SAFE_TOP, AVAIL_W, Inches(0.8)).text_frame
            set_p_format(p_title.paragraphs[0], title.upper(), get_dynamic_pt(title, 32), True, THEME["title"])
            
            content_top, avail_h = Inches(1.8), prs.slide_height - Inches(2.8)
            
            curr_x = SAFE_LEFT
            for st_code in p_stocks:
                add_editable_stock_tag(slide, curr_x, content_top, str(st_code), 'up')
                curr_x += Inches(1.2)
            for st_code in n_stocks:
                add_editable_stock_tag(slide, curr_x, content_top, str(st_code), 'down')
                curr_x += Inches(1.2)
                
            if p_stocks or n_stocks: 
                content_top += Inches(0.6) 
                avail_h -= Inches(0.6)

            if s_type == 'grid':
                gw, gh = AVAIL_W / 2 - Inches(0.2), avail_h / 2 - Inches(0.2)
                pos = [(0,0), (1,0), (0,1), (1,1)]
                for idx, b in enumerate(bullets[:4]):
                    gx = SAFE_LEFT + pos[idx][0] * (gw + Inches(0.4))
                    gy = content_top + pos[idx][1] * (gh + Inches(0.4))
                    draw_editable_ocbs_grid(slide, gx, gy, gw, gh, idx+1, b)

            elif s_type == 'table':
                t_data = s_info.get('table_data', [])
                nodes = t_data[1:] if len(t_data) > 1 else t_data
                if not nodes: 
                    s_type = 'text'
                else:
                    step = AVAIL_W / max(len(nodes), 1)
                    axis_y = content_top + Inches(1.0)
                    
                    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, SAFE_LEFT, axis_y, AVAIL_W, Inches(0.05))
                    line.fill.solid(); line.fill.fore_color.rgb = THEME["accent"]
                    
                    for n_idx, node in enumerate(nodes):
                        cx = SAFE_LEFT + (n_idx * step) + (step / 2)
                        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.12), axis_y - Inches(0.1), Inches(0.24), Inches(0.24))
                        dot.fill.solid(); dot.fill.fore_color.rgb = THEME["accent"]
                        
                        label_w = step - Inches(0.1)
                        p_yr_tf = slide.shapes.add_textbox(cx - label_w/2, axis_y - Inches(1.2), label_w, Inches(1.0)).text_frame
                        p_yr_tf.word_wrap = True
                        set_p_format(p_yr_tf.paragraphs[0], str(node[0]), get_dynamic_pt(node[0], 14), True, THEME["title"], PP_ALIGN.CENTER)
                        
                        box_w = step - Inches(0.3)
                        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx - box_w/2, axis_y + Inches(0.3), box_w, avail_h - Inches(1.4))
                        box.fill.solid(); box.fill.fore_color.rgb = THEME["card_bg"]; box.line.color.rgb = THEME["accent"]
                        
                        tf = box.text_frame
                        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.1)
                        content_txt = str(node[1]) if len(node)>1 else ""
                        set_p_format(tf.paragraphs[0], content_txt, get_dynamic_pt(content_txt, 14), None, THEME["text"])

            elif s_type == 'chart':
                c_data_dict = s_info.get('chart_data', {})
                try:
                    clean = {str(k): float(str(v).replace('%','').replace(',','')) for k, v in c_data_dict.items()}
                    cw, tw = AVAIL_W * 0.45, AVAIL_W * 0.5
                    c_data = CategoryChartData()
                    c_data.categories = list(clean.keys())
                    c_data.add_series('Value', list(clean.values()))
                    slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, SAFE_LEFT, content_top, cw, avail_h, c_data)
                except: s_type = 'text'

            if s_type == 'text':
                card_w = AVAIL_W / 2 - Inches(0.2)
                domain = str(s_info.get('company_domain', '')).strip()
                img_path = f"temp_{i}.png"; has_img = False
                if domain and len(domain) > 3 and download_company_logo(domain, img_path): has_img = True
                elif get_image_robust(str(s_info.get('image_prompt', 'finance')), "finance", img_path): has_img = True
                
                img_left, text_left = (SAFE_LEFT, SAFE_LEFT + card_w + Inches(0.4)) if i % 2 == 0 else (SAFE_LEFT + card_w + Inches(0.4), SAFE_LEFT)
                if has_img:
                    try: slide.shapes.add_picture(img_path, img_left, content_top, width=card_w); os.remove(img_path)
                    except: card_w, text_left = AVAIL_W, SAFE_LEFT
                else: card_w, text_left = AVAIL_W, SAFE_LEFT

                card_h = avail_h / max(len(bullets), 1)
                for idx, b in enumerate(bullets):
                    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, text_left, content_top + idx * card_h, card_w, card_h - Inches(0.15))
                    shape.fill.solid(); shape.fill.fore_color.rgb = THEME["card_bg"]; shape.line.color.rgb = THEME["card_border"]
                    
                    tf = shape.text_frame
                    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.1)
                    set_p_format(tf.paragraphs[0], b, get_dynamic_pt(b, 16), None, THEME["text"])

            if takeaway:
                ban = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, SAFE_LEFT, prs.slide_height - Inches(0.8), AVAIL_W, Inches(0.5))
                ban.fill.solid(); ban.fill.fore_color.rgb = THEME["card_border"]; ban.line.fill.background = None
                
                tf = ban.text_frame
                tf.margin_left = tf.margin_right = Inches(0.1)
                set_p_format(tf.paragraphs[0], f"TAKEAWAY: {takeaway}", get_dynamic_pt(takeaway, 16), True, THEME["title"], PP_ALIGN.CENTER)

        except Exception as e: continue

    out = io.BytesIO(); prs.save(out); out.seek(0)
    return out

# ==========================================
# 6. KHỞI CHẠY GIAO DIỆN
# ==========================================
def main():
    st.title("OCBS AI Report - Master Edition")
    with st.sidebar: 
        tpl = st.text_input("Template Path", "templates/template_cong_ty_moi.pptx")
        st.info("Đảm bảo đường dẫn template trỏ đến một file .pptx hợp lệ.")
        
    uf = st.file_uploader("Upload Docx", type="docx")
    if uf and st.button("Generate", type="primary"):
        with st.status("Processing...") as status:
            try:
                js = get_slide_json_from_llama3(uf, status)
                if js and len(js.get('slides', [])) > 0:
                    buf = render_pptx_clean(js, tpl, status)
                    status.update(label="Done!", state="complete")
                    st.download_button("Download PPTX", data=buf, file_name="OCBS_Report_Master.pptx")
                else: status.update(label="Không thể trích xuất dữ liệu AI.", state="error")
            except Exception as e: status.update(label=f"System Error: {str(e)}", state="error")

if __name__ == "__main__":
    main()