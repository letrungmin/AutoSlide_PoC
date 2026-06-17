import streamlit as st
import docx
import json
import requests
import time
import os
import io
import random
from dotenv import load_dotenv
from openai import OpenAI
from pptx import Presentation

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# ==========================================
# 1. CẤU HÌNH BẢO MẬT & GIAO DIỆN
# ==========================================
st.set_page_config(page_title="OCBS AutoSlide AI", page_icon="📈", layout="wide")

load_dotenv()
LLAMA3_API_KEY = os.getenv("GROQ_API_KEY")
PEXELS_API_KEY = os.getenv("PEXELS_API_KEY")
LLAMA3_BASE_URL = "https://api.groq.com/openai/v1" 

THEME = {
    "bg": RGBColor(255, 255, 255),        
    "title": RGBColor(0, 51, 102),        
    "card_bg": RGBColor(248, 250, 252),   
    "card_border": RGBColor(212, 175, 55),
    "text": RGBColor(51, 65, 85),         
    "accent": RGBColor(0, 153, 117),      
    "chart": [RGBColor(0, 153, 117), RGBColor(212, 175, 55), RGBColor(0, 51, 102), RGBColor(148, 163, 184)]
}
client = OpenAI(api_key=LLAMA3_API_KEY, base_url=LLAMA3_BASE_URL)

def download_image(keyword, filepath):
    try:
        # Làm nhẹ đuôi tìm kiếm và chọn ngẫu nhiên để không bị trùng ảnh
        search_query = f"{keyword} finance business"
        res = requests.get(f"https://api.pexels.com/v1/search?query={search_query}&per_page=15&orientation=landscape", headers={"Authorization": PEXELS_API_KEY})
        data = res.json()
        photos = data.get("photos", [])
        if photos:
            photo = random.choice(photos) # Thuật toán xáo trộn ảnh
            img_data = requests.get(photo["src"]["large"]).content
            with open(filepath, 'wb') as handler:
                handler.write(img_data)
            return True
    except: pass
    return False

# ==========================================
# 2. KIẾN TRÚC AI: LÀM SẠCH VÀ ÉP KHUÔN
# ==========================================
def get_semantic_chunks_from_docx(file_stream, max_words=600):
    doc = docx.Document(file_stream)
    chunks, current_chunk, current_word_count = [], [], 0
    for para in doc.paragraphs:
        text = para.text.strip()
        if text: 
            words = text.split()
            # [FIX LỖI DÍNH CHỮ]: Bơm thêm khoảng trắng vào cuối câu
            current_chunk.append(text + " ") 
            current_word_count += len(words)
            if current_word_count >= max_words:
                chunks.append("\n\n".join(current_chunk))
                current_chunk, current_word_count = [], 0
    if current_chunk: chunks.append("\n\n".join(current_chunk))
    return chunks

def get_slide_json_from_llama3(file_stream, status_container):
    chunks = get_semantic_chunks_from_docx(file_stream)
    all_slides = []
    
    system_prompt = """
    Bạn là Giám đốc Chiến lược. Chuyển văn bản thành JSON Presentation.
    
    KỶ LUẬT THÉP:
    1. 1 KHỐI DỮ LIỆU = 1 SLIDE DUY NHẤT. Bắt buộc nén thông tin.
    2. 'bullets': LUÔN PHẢI CÓ để làm dự phòng (kể cả slide chart/table). Viết 2-3 ý siêu ngắn.
    3. 'takeaway': Tự suy luận 1 câu chốt hạ chiến lược. TUYỆT ĐỐI KHÔNG COPY BẤT KỲ CÂU NÀO TRONG LỜI NHẮC NÀY.
    4. NHÃN BIỂU ĐỒ (chart_data): Các key PHẢI cực ngắn (1-3 từ).
    
    PHÂN LOẠI BỐ CỤC:
    - Có năm (2026, 2027) -> "type": "table"
    - Có Tỷ lệ % (45%, 50%) -> "type": "chart"
    - Còn lại -> "type": "text"
    
    JSON MẪU BẮT BUỘC:
    {
        "slides": [
            {
                "title": "Bối cảnh", 
                "type": "text", 
                "image_keyword": "market analysis",
                "bullets": ["Ý 1", "Ý 2"],
                "takeaway": "Chốt hạ chiến lược sâu sắc",
                "speaker_notes": "Nguyên văn gốc"
            },
            {
                "title": "Lộ trình", 
                "type": "table", 
                "table_data": [["Năm", "Nội dung"], ["2026", "Làm A"], ["2027", "Làm B"]],
                "bullets": ["Ý 1", "Ý 2"],
                "takeaway": "Chốt hạ",
                "speaker_notes": "Nguyên văn gốc"
            },
            {
                "title": "Ngân sách", 
                "type": "chart", 
                "chart_data": {"Hạ tầng": 40, "Bảo mật": 60},
                "bullets": ["Ý 1", "Ý 2"],
                "takeaway": "Chốt hạ",
                "speaker_notes": "Nguyên văn gốc"
            }
        ]
    }
    """
    
    progress_bar = status_container.progress(0)
    for i, chunk in enumerate(chunks):
        status_container.write(f"🧠 Đang định tuyến Đa Bố Cục (Gói {i+1}/{len(chunks)})...")
        retries = 3
        while retries > 0:
            try:
                res = client.chat.completions.create(
                    model="llama-3.1-8b-instant", 
                    messages=[
                        {"role": "system", "content": system_prompt}, 
                        {"role": "user", "content": f"Khối dữ liệu:\n{chunk}"}
                    ],
                    temperature=0.3, # Mở khóa sáng tạo để Takeaway mượt hơn
                    response_format={"type": "json_object"}
                )
                data = json.loads(res.choices[0].message.content)
                if "slides" in data: all_slides.extend(data["slides"])
                
                progress_bar.progress((i + 1) / len(chunks))
                if i < len(chunks) - 1: time.sleep(3) 
                break
            except Exception as e:
                err_msg = str(e)
                if "429" in err_msg or "413" in err_msg:
                    status_container.warning(f"⚠️ Tránh Rate Limit, ngủ đông 15s... (Còn {retries-1} lần thử)")
                    time.sleep(15)
                    retries -= 1
                else: break
    return all_slides if all_slides else None

# ==========================================
# 3. ENGINE RENDER: TỌA ĐỘ ĐỘNG (DYNAMIC LAYOUT)
# ==========================================
def render_pptx_clean(slide_data, template_path, status_container):
    status_container.write("🎨 Đang kết xuất Giao diện Pixel-Perfect...")
    prs = Presentation(template_path)
    slide_w, slide_h = prs.slide_width, prs.slide_height
    
    # [FIX LỖI LỆCH PHẢI]: Canh lề chuẩn 0.8 inches
    SAFE_LEFT = Inches(0.8) 
    SAFE_TOP = Inches(1.2) 
    AVAIL_W = slide_w - Inches(1.6) # Căn giữa hoàn hảo
    
    try:
        title_slide = prs.slides.add_slide(prs.slide_layouts[6]) 
        p_title = title_slide.shapes.add_textbox(SAFE_LEFT, Inches(3.0), AVAIL_W, Inches(2.0)).text_frame.paragraphs[0]
        p_title.text = "CHIẾN LƯỢC CHUYỂN ĐỔI SỐ OCBS"
        p_title.font.size, p_title.font.bold, p_title.font.color.rgb = Pt(44), True, THEME["title"]
        p_title.alignment = PP_ALIGN.CENTER # Căn giữa trang bìa
    except: pass

    slides_list = slide_data.get('slides', slide_data) if isinstance(slide_data, dict) else slide_data

    for i, slide_info in enumerate(slides_list):
        if not isinstance(slide_info, dict): continue
        
        slide_type = slide_info.get('type', 'text')
        title_text = slide_info.get('title', 'Untitled')
        bullets = slide_info.get('bullets', [])
        takeaway = slide_info.get('takeaway', '')
        speaker_notes = slide_info.get('speaker_notes', '') 

        if not bullets and speaker_notes:
            bullets = [speaker_notes.split('.')[0] + "."] 
        
        if slide_type == 'chart' and not slide_info.get('chart_data'): slide_type = 'text'
        elif slide_type == 'table' and len(slide_info.get('table_data', [])) <= 1: slide_type = 'text'

        try: slide = prs.slides.add_slide(prs.slide_layouts[6])
        except: slide = prs.slides.add_slide(prs.slide_layouts[1])
            
        if speaker_notes:
            slide.notes_slide.notes_text_frame.text = f"--- NỘI DUNG NGUYÊN BẢN ---\n{speaker_notes}"

        p_title = slide.shapes.add_textbox(SAFE_LEFT, SAFE_TOP, AVAIL_W, Inches(0.8)).text_frame.paragraphs[0]
        p_title.text, p_title.font.size, p_title.font.bold, p_title.font.color.rgb = title_text.upper(), Pt(28), True, THEME["title"]
        
        content_top = Inches(2.2) 
        avail_h = slide_h - content_top - Inches(1.2)
        
        # --- BỐ CỤC 1: TEXT ---
        if slide_type == 'text':
            if not bullets: continue
            gap = Inches(0.2)
            card_w = (AVAIL_W / 2) - Inches(0.2)
            
            img_path = f"temp_img_{int(time.time())}_{i}.jpg"
            if download_image(slide_info.get('image_keyword', 'corporate success'), img_path):
                try:
                    slide.shapes.add_picture(img_path, SAFE_LEFT + card_w + Inches(0.4), content_top, width=card_w, height=avail_h)
                    os.remove(img_path) 
                except: card_w = AVAIL_W 
            else: card_w = AVAIL_W

            card_h = (avail_h - (gap * (len(bullets) - 1))) / len(bullets)
            for b_idx, bullet_text in enumerate(bullets):
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, SAFE_LEFT, content_top + b_idx * (card_h + gap), card_w, card_h)
                shape.fill.solid(); shape.fill.fore_color.rgb = THEME["card_bg"]
                shape.line.color.rgb, shape.line.width = THEME["card_border"], Pt(1.5)
                
                tf = shape.text_frame
                tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.2)
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text, p.font.size, p.font.color.rgb = bullet_text.strip(), Pt(16), THEME["text"]

        # --- BỐ CỤC 2: TIMELINE ĐỘNG ---
        elif slide_type == 'table':
            t_data = slide_info.get('table_data', [])
            nodes = t_data[1:] 
            step = AVAIL_W / len(nodes)
            axis_top = Inches(3.5)
            
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, SAFE_LEFT, axis_top, AVAIL_W, Inches(0.05))
            line.fill.solid(); line.fill.fore_color.rgb = line.line.color.rgb = THEME["accent"]
            
            for n_idx, node in enumerate(nodes):
                cx = SAFE_LEFT + (n_idx * step) + (step / 2)
                dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.12), axis_top - Inches(0.10), Inches(0.24), Inches(0.24))
                dot.fill.solid(); dot.fill.fore_color.rgb = dot.line.color.rgb = THEME["accent"]
                
                p_yr = slide.shapes.add_textbox(cx - Inches(1.0), axis_top - Inches(0.8), Inches(2.0), Inches(0.6)).text_frame.paragraphs[0]
                p_yr.text, p_yr.alignment, p_yr.font.size, p_yr.font.bold, p_yr.font.color.rgb = str(node[0]), PP_ALIGN.CENTER, Pt(24), True, THEME["title"]
                
                content_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx - ((step - Inches(0.2))/2), axis_top + Inches(0.3), step - Inches(0.2), Inches(2.2))
                content_box.fill.solid(); content_box.fill.fore_color.rgb, content_box.line.color.rgb = THEME["card_bg"], THEME["accent"]
                
                p_c = content_box.text_frame.paragraphs[0]
                p_c.text, p_c.font.size, p_c.font.color.rgb = str(node[1]) if len(node) > 1 else "", Pt(14), THEME["text"]
                content_box.text_frame.word_wrap = True

        # --- BỐ CỤC 3: CHART TỶ LỆ VÀNG (CHỐNG OVERLAP) ---
        elif slide_type == 'chart':
            raw_chart_data = slide_info.get('chart_data', {})
            chart_data_dict = {}
            if isinstance(raw_chart_data, dict): chart_data_dict = raw_chart_data
            elif isinstance(raw_chart_data, list):
                for item in raw_chart_data:
                    if isinstance(item, dict): chart_data_dict.update(item)
            
            if chart_data_dict:
                try:
                    clean_dict = {str(k): float(str(v).replace('%', '').replace(',', '').strip()) for k, v in chart_data_dict.items()}
                    max_label, max_val = max(clean_dict.items(), key=lambda x: x[1])
                    display_val = int(max_val) if max_val.is_integer() else max_val
                    
                    # [FIX LỖI CHỒNG CHÉO]: Tính theo tỷ lệ % chiều ngang màn hình
                    kpi_w = AVAIL_W * 0.25
                    chart_w = AVAIL_W * 0.35
                    cmt_w = AVAIL_W * 0.35
                    
                    # 1. Cột KPI bên trái
                    tf_kpi = slide.shapes.add_textbox(SAFE_LEFT, content_top, kpi_w, avail_h).text_frame
                    p_kpi_num = tf_kpi.paragraphs[0]
                    p_kpi_num.text, p_kpi_num.font.size, p_kpi_num.font.bold, p_kpi_num.font.color.rgb = f"{display_val}%", Pt(72), True, THEME["accent"]
                    p_lbl = tf_kpi.add_paragraph()
                    p_lbl.text, p_lbl.font.size, p_lbl.font.color.rgb = f"{max_label.upper()}", Pt(16), THEME["text"]
                    
                    # 2. Cột Chart ở giữa
                    c_data = CategoryChartData()
                    c_data.categories, _ = list(clean_dict.keys()), c_data.add_series('Tỷ trọng', list(clean_dict.values()))
                    chart = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, SAFE_LEFT + kpi_w, content_top, chart_w, chart_w, c_data).chart
                    chart.has_legend = False
                    for idx, pt in enumerate(chart.series[0].points): pt.format.fill.solid(); pt.format.fill.fore_color.rgb = THEME["chart"][idx % len(THEME["chart"])]
                    
                    # 3. Cột Comments bên phải
                    if bullets:
                        cmt_box = slide.shapes.add_textbox(SAFE_LEFT + kpi_w + chart_w + Inches(0.2), content_top, cmt_w, avail_h).text_frame
                        cmt_box.word_wrap = True
                        p_cmt = cmt_box.paragraphs[0]
                        p_cmt.text, p_cmt.font.size, p_cmt.font.color.rgb = "\n\n".join(bullets), Pt(16), THEME["text"]
                except Exception as e: pass

        # --- DẢI CHỐT HẠ ---
        if takeaway:
            banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, SAFE_LEFT, slide_h - Inches(0.8), AVAIL_W, Inches(0.5))
            banner.fill.solid(); banner.fill.fore_color.rgb, banner.line.fill.background = THEME["card_border"], None
            p_ban = banner.text_frame.paragraphs[0]
            banner.text_frame.margin_left = Inches(0.2)
            display_takeaway = takeaway if "Tự viết" not in takeaway else "Ứng dụng công nghệ để mở rộng lợi thế cạnh tranh cốt lõi."
            p_ban.text, p_ban.font.size, p_ban.font.bold, p_ban.font.color.rgb = f"💡 TAKEAWAY: {display_takeaway}", Pt(14), True, THEME["title"]

    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

# ==========================================
# 4. TRẠM ĐIỀU KHIỂN STREAMLIT WEB APP
# ==========================================
def main():
    st.title("🌟 AI AutoSlide Enterprise")
    st.markdown("Hệ thống chuyển đổi Báo cáo Word thành Slide PowerPoint chuẩn C-Suite. Tích hợp phân loại tự động Lộ trình (Timeline) và Biểu đồ (Chart).")

    with st.sidebar:
        st.header("⚙️ Thiết lập Hệ thống")
        template_path = st.text_input("Đường dẫn file Template (.pptx)", "/Users/trungmin/Downloads/AutoSlide_PoC/templates/template_cong_ty_moi.pptx")
        st.info("💡 Hệ thống sử dụng Llama 3.1 8B để trích xuất dữ liệu đa định dạng và Pexels API để nhúng hình ảnh chuyên nghiệp.")
        st.caption(f"Trạng thái Bảo mật: Đã nạp thành công từ .env")

    uploaded_file = st.file_uploader("Kéo thả file bản thảo Word (.docx) vào đây", type="docx")

    if uploaded_file is not None:
        if st.button("🚀 Kích Hoạt AI Render", use_container_width=True, type="primary"):
            
            with st.status("Hệ thống đang xử lý dữ liệu...", expanded=True) as status:
                try:
                    json_data = get_slide_json_from_llama3(uploaded_file, status)
                    
                    if json_data:
                        pptx_buffer = render_pptx_clean(json_data, template_path, status)
                        status.update(label="🎉 Quá trình tạo Slide thành công rực rỡ!", state="complete", expanded=False)
                        
                        st.success("Tác phẩm của bạn đã hoàn thiện! Lỗi tỷ lệ và hình ảnh đã được khắc phục hoàn toàn.")
                        
                        st.download_button(
                            label="⬇️ Tải Xuống File PowerPoint (.pptx)",
                            data=pptx_buffer,
                            file_name="OCBS_Bao_Cao_Pixel_Perfect.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                    else:
                        status.update(label="❌ Trích xuất dữ liệu thất bại. AI không thể đọc cấu trúc file.", state="error")
                except Exception as e:
                    status.update(label=f"❌ Xảy ra lỗi hệ thống: {str(e)}", state="error")

if __name__ == "__main__":
    main()