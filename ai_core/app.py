import streamlit as st
import docx
import json
import requests
import os
import io
import urllib.parse
import random
import PyPDF2
from bs4 import BeautifulSoup
from dotenv import load_dotenv
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

st.set_page_config(page_title="Universal AI Presentation Generator", layout="wide")

load_dotenv()

if "GROQ_API_KEY" in st.secrets:
    LLAMA3_API_KEY = st.secrets["GROQ_API_KEY"]
    PEXELS_API_KEY = st.secrets.get("PEXELS_API_KEY", "")
else:
    LLAMA3_API_KEY = os.getenv("GROQ_API_KEY")
    PEXELS_API_KEY = os.getenv("PEXELS_API_KEY")

if not LLAMA3_API_KEY:
    st.error("API Key is missing. Please configure GROQ_API_KEY in Streamlit Secrets.")
    st.stop()

LLAMA3_BASE_URL = "https://api.groq.com/openai/v1"
client = OpenAI(api_key=LLAMA3_API_KEY, base_url=LLAMA3_BASE_URL)

# ──────────────────────────────────────────────────────────────────────────────
# HELPERS
# ──────────────────────────────────────────────────────────────────────────────

def hex_to_rgb_color(hex_str):
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

def download_company_logo(domain, filepath):
    try:
        res = requests.get(f"https://logo.clearbit.com/{domain}?size=800", timeout=10)
        if res.status_code == 200:
            with open(filepath, 'wb') as f:
                f.write(res.content)
            return True
    except:
        pass
    return False

def generate_ai_image(image_prompt, filepath):
    try:
        style = "modern corporate finance, high-tech glass building, glowing digital data charts, professional lighting, cinematic, hyper-realistic, 8k resolution, clean minimalist design"
        encoded = urllib.parse.quote(f"{image_prompt}, {style}")
        headers = {"User-Agent": "Mozilla/5.0"}
        res = requests.get(
            f"https://image.pollinations.ai/prompt/{encoded}?width=1024&height=576&nologo=true&seed={random.randint(1, 100000)}",
            headers=headers, timeout=20
        )
        if res.status_code == 200 and len(res.content) > 5000:
            with open(filepath, 'wb') as f:
                f.write(res.content)
            return True
    except:
        pass
    return False

def download_image_pexels(keyword, filepath):
    try:
        res = requests.get(
            f"https://api.pexels.com/v1/search?query={keyword} finance business&per_page=10&orientation=landscape",
            headers={"Authorization": PEXELS_API_KEY}
        )
        photos = res.json().get("photos", [])
        if photos:
            img_data = requests.get(random.choice(photos)["src"]["large"]).content
            with open(filepath, 'wb') as f:
                f.write(img_data)
            return True
    except:
        pass
    return False

def get_image_robust(prompt, keyword, filepath):
    if generate_ai_image(prompt, filepath):
        return True
    return download_image_pexels(keyword, filepath)

# ──────────────────────────────────────────────────────────────────────────────
# TEXT EXTRACTION – với overlap để giữ context giữa các chunk
# ──────────────────────────────────────────────────────────────────────────────

def chunk_paragraphs(paragraphs, max_words=600, overlap_paras=2):
    """
    FIX: Chunk theo đơn vị đoạn văn (không phải từ), có overlap để giữ context.
    overlap_paras: số đoạn cuối chunk trước được lặp lại ở đầu chunk sau.
    """
    chunks = []
    current_paras = []
    current_len = 0

    for para in paragraphs:
        words = para.split()
        if current_len + len(words) > max_words and current_paras:
            chunks.append("\n".join(current_paras))
            # Overlap: giữ lại overlap_paras đoạn cuối làm đầu chunk mới
            current_paras = current_paras[-overlap_paras:]
            current_len = sum(len(p.split()) for p in current_paras)
        current_paras.append(para)
        current_len += len(words)

    if current_paras:
        chunks.append("\n".join(current_paras))

    return chunks if chunks else [""]

def get_semantic_chunks_from_docx(file_stream, max_words=600):
    file_stream.seek(0)
    doc = docx.Document(file_stream)
    paragraphs = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
    return chunk_paragraphs(paragraphs, max_words)

def get_semantic_chunks_from_pdf(file_stream, max_words=600):
    file_stream.seek(0)
    reader = PyPDF2.PdfReader(file_stream)
    paragraphs = []
    for page in reader.pages:
        extracted = page.extract_text()
        if extracted:
            for line in extracted.split('\n'):
                line = line.strip()
                if line:
                    paragraphs.append(line)
    return chunk_paragraphs(paragraphs, max_words)

def get_semantic_chunks_from_url(url, max_words=600):
    headers = {"User-Agent": "Mozilla/5.0"}
    res = requests.get(url, headers=headers, timeout=15)
    res.raise_for_status()
    soup = BeautifulSoup(res.content, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.extract()
    paragraphs = [t.strip() for t in soup.get_text(separator='\n').split('\n') if t.strip()]
    return chunk_paragraphs(paragraphs, max_words)

# ──────────────────────────────────────────────────────────────────────────────
# AI – PROMPT CẢI TIẾN: kiểm soát độ dài bullet và số lượng slide
# ──────────────────────────────────────────────────────────────────────────────

def get_slide_json_from_llama3(chunks, status_container):
    all_slides = []
    system_prompt = """
You extract structured presentation data from documents. Output ONLY valid JSON, no markdown fences.

OUTPUT FORMAT:
{
  "slides": [
    {
      "title": "Slide title (max 60 chars)",
      "type": "text",
      "company_domain": "",
      "image_prompt": "10 descriptive english words for stock photo",
      "takeaway": "One concise concluding sentence (max 120 chars)",
      "positive_stocks": [],
      "negative_stocks": [],
      "bullets": [
        "Complete idea — between 30 and 150 characters each"
      ],
      "table_data": [],
      "chart_data": {}
    }
  ]
}

STRICT RULES:
1. LANGUAGE: Preserve ALL input language and diacritics exactly (Vietnamese, etc.)
2. BULLETS: 3 to 5 bullets per slide. Each bullet MUST be 30-150 characters — a complete idea.
   Never leave bullets empty []. Never cut a sentence mid-way.
3. TYPE:
   - Timeline / roadmap / stages → "table", table_data: [["Label", "Description"], ...]
   - Numeric proportions / percentages → "chart", chart_data: {"Label": number, ...}
   - Everything else → "text"
4. SPLIT large topics across multiple slides (one clear theme per slide).
5. title: max 60 characters. takeaway: max 120 characters.
6. image_prompt: 10 English words describing a relevant business/tech photo.
"""

    for i, chunk in enumerate(chunks):
        status_container.write(f"🔍 Phân tích nội dung (Batch {i+1}/{len(chunks)})...")
        retries = 3
        while retries > 0:
            try:
                res = client.chat.completions.create(
                    model="llama-3.1-8b-instant",
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": chunk}
                    ],
                    temperature=0.15,
                    response_format={"type": "json_object"}
                )
                content = res.choices[0].message.content.strip()
                content = content.replace("```json", "").replace("```", "").strip()
                chunk_json = json.loads(content)

                if isinstance(chunk_json, dict):
                    if "slides" in chunk_json:
                        all_slides.extend(chunk_json["slides"])
                    elif "slide" in chunk_json:
                        all_slides.extend(chunk_json["slide"])
                    elif "title" in chunk_json:
                        all_slides.append(chunk_json)
                elif isinstance(chunk_json, list):
                    all_slides.extend(chunk_json)

                break
            except Exception as e:
                retries -= 1
                if retries == 0:
                    status_container.error(f"Lỗi batch {i+1}: {e}")

    return {"slides": all_slides}

# ──────────────────────────────────────────────────────────────────────────────
# PPTX RENDERING
# ──────────────────────────────────────────────────────────────────────────────

def clear_placeholders(slide):
    """Xóa tất cả placeholder trong slide (giữ nguyên background từ layout)."""
    for shape in list(slide.placeholders):
        sp = shape.element
        sp.getparent().remove(sp)

def get_dynamic_pt(text, default_size):
    """
    FIX: Kích thước font bảo thủ hơn để tránh tràn.
    """
    length = len(str(text))
    if length < 35:
        return Pt(default_size)
    elif length < 90:
        return Pt(max(13, default_size - 2))
    elif length < 160:
        return Pt(max(12, default_size - 4))
    else:
        return Pt(max(11, default_size - 5))

def set_p_format(paragraph, text, font_size, bold=False, color_rgb=None, alignment=None):
    """
    FIX BUG #1 – FONT: Áp dụng font ở CẢ cấp paragraph VÀ run.
    
    Vấn đề cũ: paragraph.font.name = 'Arial' chỉ set cấp paragraph.
    Nhưng paragraph.runs[0].font.name = None (kế thừa), LibreOffice/PowerPoint
    đôi khi không đọc font kế thừa đúng → chữ hiển thị sai font.
    
    Fix: đặt font trực tiếp trên run (ưu tiên cao hơn paragraph-level).
    """
    paragraph.text = str(text)

    # Paragraph-level (fallback/inheritance)
    paragraph.font.name = 'Arial'
    paragraph.font.size = font_size
    paragraph.font.bold = bold
    if color_rgb:
        paragraph.font.color.rgb = color_rgb

    # Run-level
    for run in paragraph.runs:
        run.font.name = 'Arial'
        run.font.size = font_size
        run.font.bold = bold
        if color_rgb:
            run.font.color.rgb = color_rgb

    if alignment:
        paragraph.alignment = alignment


def _configure_tf(tf, word_wrap=True):
    """
    FIX BUG #2 – WORD WRAP: Set word_wrap trên TextFrame object (không phải XML element).
    
    Vấn đề cũ: paragraph._parent.word_wrap = True không hoạt động vì
    paragraph._parent là XML element (CT_TxBody), không phải TextFrame Python object.
    TextFrame.word_wrap là Python property, không phải XML attribute.
    
    Fix: gọi tf.word_wrap = True trực tiếp trên TextFrame khi tạo shape.
    """
    tf.word_wrap = word_wrap
    # Không dùng TEXT_TO_FIT_SHAPE (không tin cậy khi render bên ngoài PowerPoint)
    # Thay vào đó dùng sizing bảo thủ + word_wrap


def add_editable_stock_tag(slide, x, y, stock_code, trend, theme):
    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.1), Inches(0.4))
    bg_color = theme["accent"] if trend == 'up' else RGBColor(220, 38, 38)
    tag.fill.solid()
    tag.fill.fore_color.rgb = bg_color
    tag.line.fill.background = None
    tf = tag.text_frame
    _configure_tf(tf)
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    set_p_format(tf.paragraphs[0], stock_code.upper(), Pt(11), True, RGBColor(255, 255, 255), PP_ALIGN.CENTER)


def render_pptx_clean(slide_data, template_source, report_title, theme):
    prs = Presentation(template_source)
    SLIDE_W = prs.slide_width
    SLIDE_H = prs.slide_height
    SAFE_LEFT = Inches(0.5)
    SAFE_TOP  = Inches(0.75)
    AVAIL_W   = SLIDE_W - Inches(1.0)

    # ──────────────────────────────────────────────────────────────
    # FIX BUG #3 – SLIDE ĐẦU TRỐNG: Dùng slide có sẵn của template
    # làm cover thay vì add_slide() (sẽ để lại slide đầu trống).
    # ──────────────────────────────────────────────────────────────
    if len(prs.slides) > 0:
        cover_slide = prs.slides[0]
        clear_placeholders(cover_slide)
    else:
        cover_slide = prs.slides.add_slide(prs.slide_layouts[6])
        clear_placeholders(cover_slide)

    cover_tb = cover_slide.shapes.add_textbox(SAFE_LEFT, Inches(2.8), AVAIL_W, Inches(2.2))
    cover_tf = cover_tb.text_frame
    _configure_tf(cover_tf)
    cover_tf.margin_left = cover_tf.margin_right = Inches(0.2)
    set_p_format(cover_tf.paragraphs[0], report_title.upper(), Pt(48), True, theme["title"], PP_ALIGN.CENTER)

    # ──────────────────────────────────────────────────────────────
    # Vùng layout cố định cho tất cả content slides
    # ──────────────────────────────────────────────────────────────
    TITLE_H     = Inches(0.85)       # chiều cao textbox tiêu đề
    TITLE_TOP   = SAFE_TOP           # 0.75"
    CONT_TOP_BASE = TITLE_TOP + TITLE_H + Inches(0.15)  # ~1.75"
    TAKEAWAY_H  = Inches(0.48)
    TAKEAWAY_Y  = SLIDE_H - TAKEAWAY_H - Inches(0.12)   # 6.9" from top
    CONT_AVAIL_H = TAKEAWAY_Y - CONT_TOP_BASE - Inches(0.1)  # ~5.0"

    slides_list = slide_data.get('slides', []) if isinstance(slide_data, dict) else slide_data

    for i, s_info in enumerate(slides_list):
        if not isinstance(s_info, dict):
            continue
        try:
            s_type = str(s_info.get('type') or 'text')
            if s_type not in ['table', 'chart']:
                s_type = 'text'

            title    = str(s_info.get('title') or 'Content Slide')
            takeaway = str(s_info.get('takeaway') or '').strip()

            # ── Xử lý bullets: làm sạch, ghép đoạn ngắn, giới hạn 5 ──
            raw_bullets = [str(b).strip() for b in (s_info.get('bullets') or []) if str(b).strip()]
            cleaned = []
            for b in raw_bullets:
                if len(b) < 25 and cleaned:
                    cleaned[-1] += ' ' + b  # ghép đoạn quá ngắn
                else:
                    cleaned.append(b)
            if not cleaned:
                cleaned = ["Nội dung đang được tổng hợp từ tài liệu gốc."]
            bullets = cleaned[:5]  # cap 5 bullets để tránh overflow

            p_stocks = s_info.get('positive_stocks') or []
            n_stocks = s_info.get('negative_stocks') or []

            # Thêm slide mới (dùng layout Blank để giữ background template)
            try:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                clear_placeholders(slide)
            except:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                clear_placeholders(slide)

            # ── Tiêu đề slide ────────────────────────────────────────
            title_tb = slide.shapes.add_textbox(SAFE_LEFT, TITLE_TOP, AVAIL_W, TITLE_H)
            title_tf = title_tb.text_frame
            _configure_tf(title_tf)
            title_tf.margin_left  = Inches(0.05)
            title_tf.margin_right = Inches(0.05)
            title_tf.margin_top   = Inches(0.04)
            title_tf.margin_bottom = Inches(0.04)
            title_font_size = get_dynamic_pt(title, 28)
            set_p_format(title_tf.paragraphs[0], title.upper(), title_font_size, True, theme["title"])

            # ── Stock tags (nếu có) ──────────────────────────────────
            cont_top  = CONT_TOP_BASE
            avail_h   = CONT_AVAIL_H

            curr_x = SAFE_LEFT
            for st_code in p_stocks:
                add_editable_stock_tag(slide, curr_x, cont_top, str(st_code), 'up', theme)
                curr_x += Inches(1.2)
            for st_code in n_stocks:
                add_editable_stock_tag(slide, curr_x, cont_top, str(st_code), 'down', theme)
                curr_x += Inches(1.2)
            if p_stocks or n_stocks:
                cont_top += Inches(0.5)
                avail_h  -= Inches(0.5)

            # ── TABLE SLIDE ──────────────────────────────────────────
            if s_type == 'table':
                t_data = s_info.get('table_data', [])
                nodes = t_data[1:] if len(t_data) > 1 else t_data
                if not nodes:
                    s_type = 'text'
                else:
                    n_nodes = min(len(nodes), 6)  # tối đa 6 cột
                    nodes = nodes[:n_nodes]
                    step = AVAIL_W / n_nodes
                    axis_y = cont_top + Inches(1.0)

                    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, SAFE_LEFT, axis_y, AVAIL_W, Inches(0.05))
                    line.fill.solid()
                    line.fill.fore_color.rgb = theme["accent"]

                    for n_idx, node in enumerate(nodes):
                        cx = SAFE_LEFT + (n_idx * step) + (step / 2)

                        dot = slide.shapes.add_shape(
                            MSO_SHAPE.OVAL,
                            cx - Inches(0.12), axis_y - Inches(0.12),
                            Inches(0.25), Inches(0.25)
                        )
                        dot.fill.solid()
                        dot.fill.fore_color.rgb = theme["accent"]
                        dot.line.fill.background = None

                        label_w = step - Inches(0.1)
                        label_tb = slide.shapes.add_textbox(
                            cx - label_w / 2, axis_y - Inches(1.1), label_w, Inches(0.95)
                        )
                        label_tf = label_tb.text_frame
                        _configure_tf(label_tf)
                        label_tf.margin_left = label_tf.margin_right = Inches(0.05)
                        set_p_format(
                            label_tf.paragraphs[0],
                            str(node[0]),
                            get_dynamic_pt(node[0], 13), True, theme["title"], PP_ALIGN.CENTER
                        )

                        box_w = step - Inches(0.25)
                        box_h = avail_h - Inches(1.5)
                        box = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE,
                            cx - box_w / 2, axis_y + Inches(0.3),
                            box_w, max(box_h, Inches(0.6))
                        )
                        box.fill.solid()
                        box.fill.fore_color.rgb = theme["card_bg"]
                        box.line.color.rgb = theme["accent"]
                        box_tf = box.text_frame
                        _configure_tf(box_tf)
                        box_tf.margin_left = box_tf.margin_right = Inches(0.1)
                        box_tf.margin_top  = box_tf.margin_bottom = Inches(0.08)
                        content_txt = str(node[1]) if len(node) > 1 else ""
                        set_p_format(
                            box_tf.paragraphs[0],
                            content_txt,
                            get_dynamic_pt(content_txt, 12), False, theme["text"]
                        )

            # ── CHART SLIDE ──────────────────────────────────────────
            elif s_type == 'chart':
                c_data_dict = s_info.get('chart_data', {})
                try:
                    clean = {
                        str(k): float(str(v).replace('%', '').replace(',', ''))
                        for k, v in c_data_dict.items()
                    }
                    cw = AVAIL_W * 0.55
                    c_data = CategoryChartData()
                    c_data.categories = list(clean.keys())
                    c_data.add_series('Value', list(clean.values()))
                    slide.shapes.add_chart(
                        XL_CHART_TYPE.BAR_CLUSTERED,
                        SAFE_LEFT, cont_top, cw, avail_h, c_data
                    )
                except:
                    s_type = 'text'

            # ── TEXT SLIDE ───────────────────────────────────────────
            if s_type == 'text':
                # Tải ảnh
                card_w = AVAIL_W / 2 - Inches(0.2)
                domain = str(s_info.get('company_domain', '')).strip()
                img_path = f"temp_{i}.png"
                has_img = False

                if domain and len(domain) > 3:
                    has_img = download_company_logo(domain, img_path)
                if not has_img:
                    has_img = get_image_robust(
                        str(s_info.get('image_prompt', 'business technology digital')),
                        "business", img_path
                    )

                if has_img:
                    img_left  = SAFE_LEFT if i % 2 == 0 else SAFE_LEFT + card_w + Inches(0.4)
                    text_left = SAFE_LEFT + card_w + Inches(0.4) if i % 2 == 0 else SAFE_LEFT
                    try:
                        # FIX: đặt cả width VÀ height để ảnh không bị crop
                        img_h = avail_h * 0.88
                        slide.shapes.add_picture(img_path, img_left, cont_top, width=card_w, height=img_h)
                        os.remove(img_path)
                    except:
                        has_img = False

                if not has_img:
                    card_w    = AVAIL_W
                    text_left = SAFE_LEFT

                # ── FIX BUG #4 – TEXT OVERFLOW ──────────────────────
                # Tính chiều cao card từ không gian thực tế, đảm bảo tối thiểu.
                n = len(bullets)
                GAP = Inches(0.1)                             # khoảng cách giữa các card
                MIN_CARD_H = Inches(0.58)                     # chiều cao tối thiểu mỗi card
                total_gap = GAP * (n - 1) if n > 1 else Emu(0)
                card_h = max(MIN_CARD_H, (avail_h - total_gap) / n)

                for idx, b in enumerate(bullets):
                    y_pos = cont_top + idx * (card_h + GAP)
                    shape = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        text_left, y_pos, card_w, card_h
                    )
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = theme["card_bg"]
                    shape.line.color.rgb = theme["card_border"]

                    tf = shape.text_frame
                    _configure_tf(tf)          # ← word_wrap = True đúng cách
                    tf.margin_left  = Inches(0.12)
                    tf.margin_right = Inches(0.12)
                    tf.margin_top   = Inches(0.08)
                    tf.margin_bottom = Inches(0.06)

                    # Font nhỏ hơn khi card thấp
                    if card_h < Inches(0.75):
                        base_size = 12
                    else:
                        base_size = 14
                    set_p_format(tf.paragraphs[0], b, get_dynamic_pt(b, base_size), False, theme["text"])

            # ── TAKEAWAY BAR ─────────────────────────────────────────
            if takeaway:
                ban = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    SAFE_LEFT, TAKEAWAY_Y, AVAIL_W, TAKEAWAY_H
                )
                ban.fill.solid()
                ban.fill.fore_color.rgb = theme["card_border"]
                ban.line.fill.background = None

                tf = ban.text_frame
                _configure_tf(tf)
                tf.margin_left  = Inches(0.15)
                tf.margin_right = Inches(0.15)
                tf.margin_top   = Inches(0.04)
                tf.margin_bottom = Inches(0.04)
                set_p_format(
                    tf.paragraphs[0],
                    f"TAKEAWAY: {takeaway}",
                    get_dynamic_pt(takeaway, 12), True, theme["title"], PP_ALIGN.CENTER
                )

        except Exception:
            continue

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out


# ──────────────────────────────────────────────────────────────────────────────
# STREAMLIT UI  (không thay đổi logic, chỉ cải tiến label)
# ──────────────────────────────────────────────────────────────────────────────

def main():
    st.title("Universal AI Presentation Generator")

    if "slide_data" not in st.session_state:
        st.session_state.slide_data = None
    if "edited_data" not in st.session_state:
        st.session_state.edited_data = None
    if "ppt_buffer" not in st.session_state:
        st.session_state.ppt_buffer = None

    with st.sidebar:
        st.header("Report Configuration")
        report_title = st.text_input("Cover Title:", "ANALYSIS REPORT")

        template_dir = "ai_core/templates"
        available_templates = {}
        if os.path.exists(template_dir):
            for file in os.listdir(template_dir):
                if file.endswith(".pptx"):
                    display_name = file.replace(".pptx", "").replace("_", " ").title()
                    available_templates[display_name] = os.path.join(template_dir, file)

        if not available_templates:
            available_templates["Default Template"] = "ai_core/templates/template_cong_ty_moi.pptx"

        selected_template_name = st.selectbox("Select System Template:", list(available_templates.keys()))
        selected_template_path = available_templates[selected_template_name]

        custom_template = st.file_uploader("Upload Custom Template (.pptx)", type="pptx")

        st.markdown("---")
        st.subheader("Brand Kit")
        col1, col2 = st.columns(2)
        with col1:
            title_hex = st.color_picker("Title", "#003366")
            text_hex  = st.color_picker("Text",  "#334155")
        with col2:
            accent_hex = st.color_picker("Accent", "#009975")
            border_hex = st.color_picker("Border", "#D4AF37")

        custom_theme = {
            "title":       hex_to_rgb_color(title_hex),
            "accent":      hex_to_rgb_color(accent_hex),
            "card_bg":     RGBColor(248, 250, 252),
            "card_border": hex_to_rgb_color(border_hex),
            "text":        hex_to_rgb_color(text_hex),
            "chart":       [
                hex_to_rgb_color(accent_hex),
                hex_to_rgb_color(border_hex),
                hex_to_rgb_color(title_hex),
                RGBColor(148, 163, 184)
            ]
        }

    st.subheader("Data Source")
    input_source = st.radio(
        "Select input method:",
        ["File Upload (DOCX/PDF)", "Website URL"],
        label_visibility="collapsed"
    )

    uf = None
    url_input = ""

    if input_source == "File Upload (DOCX/PDF)":
        uf = st.file_uploader("Upload Document", type=["docx", "pdf"])
    else:
        url_input = st.text_input("Enter Web Article URL")

    if (uf or url_input) and st.button("1. Analyze Content", type="primary"):
        with st.status("Đang phân tích nội dung...") as status:
            try:
                chunks = []
                if uf and uf.name.endswith('.docx'):
                    chunks = get_semantic_chunks_from_docx(uf)
                elif uf and uf.name.endswith('.pdf'):
                    chunks = get_semantic_chunks_from_pdf(uf)
                elif url_input:
                    status.update(label="Đang tải nội dung trang web...")
                    chunks = get_semantic_chunks_from_url(url_input)

                if not chunks or all(not c.strip() for c in chunks):
                    status.update(label="Không tìm thấy nội dung có thể đọc.", state="error")
                    st.stop()

                js = get_slide_json_from_llama3(chunks, status)
                if js and len(js.get('slides', [])) > 0:
                    st.session_state.slide_data = js
                    st.session_state.ppt_buffer = None
                    status.update(label=f"Trích xuất xong {len(js['slides'])} slide. Xem lại bên dưới.", state="complete")
                else:
                    status.update(label="Hệ thống không trích xuất được dữ liệu có cấu trúc.", state="error")
            except Exception as e:
                status.update(label=f"Lỗi hệ thống: {str(e)}", state="error")

    if st.session_state.slide_data:
        st.markdown("---")
        st.subheader("2. Review & Edit Content")

        edited_slides = []
        for i, slide in enumerate(st.session_state.slide_data['slides']):
            with st.expander(f"Slide {i+1}: {slide.get('title', 'Untitled')}", expanded=False):
                new_title    = st.text_input("Title",    slide.get('title', ''),    key=f"title_{i}")
                new_takeaway = st.text_input("Takeaway", slide.get('takeaway', ''), key=f"takeaway_{i}")

                bullets_str     = "\n".join(slide.get('bullets', []))
                new_bullets_str = st.text_area("Bullets (one per line)", bullets_str, height=120, key=f"bullets_{i}")

                new_slide = slide.copy()
                new_slide['title']    = new_title
                new_slide['takeaway'] = new_takeaway
                new_slide['bullets']  = [b.strip() for b in new_bullets_str.split("\n") if b.strip()]
                edited_slides.append(new_slide)

        st.session_state.edited_data = {"slides": edited_slides}

        if st.button("3. Generate PowerPoint"):
            template_source = custom_template if custom_template else selected_template_path
            with st.spinner("Đang render slides..."):
                try:
                    buf = render_pptx_clean(
                        st.session_state.edited_data,
                        template_source,
                        report_title,
                        custom_theme
                    )
                    st.session_state.ppt_buffer = buf
                    st.success("PowerPoint đã được tạo thành công!")
                except Exception as e:
                    st.error(f"Render Error: {str(e)}")

    if st.session_state.ppt_buffer:
        st.download_button(
            label="Download PowerPoint File",
            data=st.session_state.ppt_buffer,
            file_name="Universal_Presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            type="primary"
        )


if __name__ == "__main__":
    main()